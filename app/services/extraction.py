"""Document text extraction service."""

from pathlib import Path

import fitz  # PyMuPDF
import structlog
from pptx import Presentation

logger = structlog.get_logger()


class ExtractionError(Exception):
    """Raised when document extraction fails."""

    pass


async def extract_content(file_path: Path, mime_type: str) -> str:
    """Extract text content from uploaded document.

    Args:
        file_path: Path to the uploaded file
        mime_type: MIME type of the file

    Returns:
        Extracted text content

    Raises:
        ExtractionError: If extraction fails
        ValueError: If file type is unsupported
    """
    logger.info("extracting_content", path=str(file_path), mime_type=mime_type)

    try:
        match mime_type:
            case "application/pdf":
                content = extract_pdf(file_path)
            case "application/vnd.openxmlformats-officedocument.presentationml.presentation":
                content = extract_pptx(file_path)
            case "text/plain":
                content = file_path.read_text(encoding="utf-8")
            case _:
                raise ValueError(f"Unsupported file type: {mime_type}")

        logger.info("extraction_complete", chars=len(content))
        return content

    except ValueError:
        raise
    except Exception as e:
        logger.error("extraction_failed", error=str(e))
        raise ExtractionError(f"Failed to extract content: {str(e)}") from e


def extract_pdf(path: Path) -> str:
    """Extract text from a PDF file using PyMuPDF.

    Args:
        path: Path to the PDF file

    Returns:
        Concatenated text from all pages
    """
    doc = fitz.open(path)
    pages = []

    for page_num, page in enumerate(doc):
        text = page.get_text()
        if text.strip():
            pages.append(f"--- Page {page_num + 1} ---\n{text}")

    doc.close()
    return "\n\n".join(pages)


def extract_pptx(path: Path) -> str:
    """Extract text from a PowerPoint file.

    Args:
        path: Path to the PPTX file

    Returns:
        Concatenated text from all slides
    """
    prs = Presentation(path)
    slides = []

    for slide_num, slide in enumerate(prs.slides, 1):
        slide_texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_texts.append(shape.text)

        if slide_texts:
            slides.append(f"--- Slide {slide_num} ---\n" + "\n".join(slide_texts))

    return "\n\n".join(slides)


# Supported MIME types for validation
SUPPORTED_MIME_TYPES = {
    "application/pdf",
    "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    "text/plain",
}


def is_supported_mime_type(mime_type: str) -> bool:
    """Check if a MIME type is supported for extraction."""
    return mime_type in SUPPORTED_MIME_TYPES
